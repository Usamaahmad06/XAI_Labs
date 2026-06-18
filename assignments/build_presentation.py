"""
Build a presentation for all XAI assignments.

Usage:
    python assignments/build_presentation.py
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
ASSIGNMENTS = REPO / "assignments"
OUT = ASSIGNMENTS / "portfolio" / "XAI_Assignments_1_to_8_Report_v2.pptx"

LABS = [
    (1, "assignment_01_causal_forecast", "Causal Forecasting"),
    (2, "assignment_02_causal_anomaly", "Explainable Anomaly Detection"),
    (3, "assignment_03_asp", "ASP Rocksample Planning"),
    (4, "assignment_04_ilasp", "ILASP + ASP Planning"),
    (5, "assignment_05_activation_rate", "Activation Rate (Explain RL)"),
    (6, "assignment_06_sr_dqn", "SR-DQN (NeSy DQN)"),
    (7, "assignment_07_nesy_distributional", "H-C51 (Distributional NeSy RL)"),
    (8, "assignment_08_stl_rl", "STL-constrained RL (Pendulum)"),
]

# Rich per-lab narrative (asked → implementation → interpretation)
LAB_NARRATIVES: dict[int, dict] = {
    1: {
        "asked": [
            "Predict building cooling demand from time-series sensor data.",
            "Compare a causal (sparse, interpretable) model against a TCN baseline.",
        ],
        "implementation": [
            "Used Tigramite Prediction + LinearRegression on causal parents of cooling_demand.",
            "Selected only variables with a directed causal link to the target.",
            "Trained a dilated TCN on all 32 input channels as the black-box baseline.",
            "Evaluated both with NMAE on the held-out test split.",
        ],
        "results": [
            "Causal model: 6 predictors, NMAE = 0.096.",
            "TCN baseline: 32 channels, NMAE = 0.130.",
            "Causal model is sparser and more accurate.",
        ],
        "interpretation": [
            "Humidity, past demand, and pricing drive predictions — not all 32 channels.",
            "Interpretability does not hurt accuracy; it improves it on this dataset.",
        ],
        "images": [
            ("predictions.png", "Actual vs predicted cooling demand — causal model tracks the signal well."),
            ("feature_importance.png", "Which causal parents contribute most to the forecast."),
        ],
        "csv": "metrics_summary.csv",
    },
    2: {
        "asked": [
            "Detect Pepper robot cyber-attacks online using explainable signals.",
            "Use causal coefficient drift rather than opaque anomaly scores.",
        ],
        "implementation": [
            "Fitted PCMCI offline to learn the causal graph and structural coefficients.",
            "Monitored sliding windows of coefficients during live operation.",
            "Flagged anomalies when coefficients deviated from the learned baseline.",
        ],
        "results": [
            "Precision: 0.64 | Recall: 0.89 | F1: 0.75.",
            "56 false positives on normal data; strong recall on attacks.",
        ],
        "interpretation": [
            "High recall means most attacks are caught; precision trade-off is acceptable for safety.",
            "Per-attack CSV shows which sensor links broke for each attack type.",
        ],
        "images": [
            ("feature_importance.png", "Coefficient drift highlights which causal links changed under attack."),
        ],
        "csv": "metrics_summary.csv",
    },
    3: {
        "asked": [
            "Plan robot actions to sample all good rocks on a grid (Rocksample).",
            "Use Answer Set Programming — no ML training.",
        ],
        "implementation": [
            "Wrote Clingo rules in rocksample.lp: preconditions, effects, and constraints.",
            "Implemented #program step(t) for legal moves and state transitions.",
            "Ran incremental ASP search from (0,0) to sample rocks 0 and 1.",
        ],
        "results": [
            "Multiple valid plans saved in plans.txt.",
            "Each plan samples good rocks 0 and 1 while respecting grid constraints.",
        ],
        "interpretation": [
            "Symbolic planning gives auditable action sequences, not a black-box policy.",
            "Clingo proves plan validity via logic — explainability by construction.",
        ],
        "images": [],
        "text_artifact": ("plans.txt", "Sample valid ASP plan (first 8 lines):"),
    },
    4: {
        "asked": [
            "Learn which rocks are good from distance/guess observations (ILASP).",
            "Then plan to sample learned-good rocks using ASP.",
        ],
        "implementation": [
            "Added mode bias to ilasp_task.las for ILP learning.",
            "Encoded learned good(R) rules in rocksample_ilp.lp with ASP step(t).",
            "Part B: planning via run_incmode.py on Windows.",
        ],
        "results": [
            "Plans target rock 2 when guess(2,90) > 80 (learned goodness).",
            "5 valid plans in plans.txt; learned rules in learned_rules.txt.",
        ],
        "interpretation": [
            "Logic learns under uncertainty; ASP executes the induced policy.",
            "Full ILASP learning optional on Linux; planning verified on Windows.",
        ],
        "images": [],
        "text_artifact": ("learned_rules.txt", "Learned ILASP rules:"),
    },
    5: {
        "asked": [
            "Explain a trained MAPPO warehouse policy: which logical rules drive which actions?",
            "Measure rule activation rate per agent and action.",
        ],
        "implementation": [
            "Implemented _does_rule_activate in parallel_runner.py.",
            "Compared observation atoms to rule bodies for each MAPPO action.",
            "Ran activation-rate experiment on RWARE with checkpoint at 1M steps.",
        ],
        "results": [
            "Agent 0: forward ~4.5%, left ~8.6%, right ~9.9%.",
            "Agent 1: load ~1.0%, right ~11.3%.",
        ],
        "interpretation": [
            "Return alone does not explain behavior — activation rate links rules to actions.",
            "Different agents activate different rule subsets on the same task.",
        ],
        "images": [
            ("plot_activation_rate.png", "Activation rate per rule/action — which heuristics the policy follows."),
        ],
        "csv": "activation_rate_20260616005507.csv",
    },
    6: {
        "asked": [
            "DoorKey MiniGrid: guide DQN exploration with symbolic heuristics (SR-DQN).",
            "When exploring, bias action sampling toward logically suggested actions.",
        ],
        "implementation": [
            "Implemented _get_random_action() in SR_DQN.py.",
            "Suggested actions get weight conf_level; others share the remainder uniformly.",
            "Trained DQN vs SR-DQN on DoorKey 5×5 for 100k steps.",
        ],
        "results": [
            "SR-DQN reaches reward ~0.7 by end of training.",
            "DQN baseline stays near 0 — random exploration fails on sparse reward.",
        ],
        "interpretation": [
            "Symbolic heuristics during exploration help discover the +1 goal reward.",
            "Explainable logic shapes learning, not only the final greedy policy.",
        ],
        "images": [
            ("rewards.png", "DQN (blue) vs SR-DQN (orange) — heuristic-guided exploration wins."),
        ],
    },
    7: {
        "asked": [
            "Combine C51 (distributional DQN) with logical heuristics on return distributions.",
            "Implement the H-C51 product on PMFs for rule-suggested actions.",
        ],
        "implementation": [
            "In h_c51_product.py get_action(): multiply rule-suggested PMFs by 1 + ε·rule_pmf.",
            "Renormalize over atoms; compute Q-values from combined distribution.",
            "Trained C51 baseline vs H-C51 for 50k steps on DoorKey 5×5.",
        ],
        "results": [
            "H-C51 sustains returns ~0.5–0.8.",
            "C51 baseline drops to ~0 after early episodes.",
        ],
        "interpretation": [
            "Heuristics shift probability mass toward high-return atoms for suggested actions.",
            "Distributional RL + logic beats plain C51 under sparse DoorKey rewards.",
        ],
        "images": [
            ("c51_vs_hc51_returns.png", "C51 baseline vs H-C51 — heuristic modulation of return distributions helps learning."),
        ],
        "csv": "metrics_summary.csv",
    },
    8: {
        "asked": [
            "Pendulum swing-up with STL constraints on angular velocity and torque.",
            "Implement STL robustness-based costs in Pendulum.py step().",
            "Compare Lagrangian TRPO (explicit costs) vs reward shaping.",
        ],
        "implementation": [
            "Built τ-horizon trajectory [thetadot, torque] each step.",
            "Computed STL robustness ρ via RTAMT (eventually always |thetadot|≤0.5, |torque|≤0.3).",
            "Mapped ρ to smooth costs: cost = tanh(−β·ρ) per constraint (β=100).",
            "Integrated standard pendulum dynamics and reward r = −(θ² + 0.1·θ̇² + 0.001·τ²).",
            "Our method: dual critics + Lagrangian TRPO (main_pendulum_tuning.py).",
            "Baseline: STL terms mixed into reward (main_pendulum_stlgym.py).",
        ],
        "results": [
            "Our (STL TRPO): return improves to ~−18; total cost → 0 (constraints satisfied).",
            "Reward shaping: return plateaus ~−65; cost stays ~−180 (weak satisfaction).",
            "50 episodes, seed 47, start_safety=100.",
        ],
        "interpretation": [
            "Higher total cost (closer to 0) means STL constraints are better satisfied.",
            "Our method separates reward maximization from constraint enforcement via dual critics.",
            "Reward shaping conflates objectives — agent optimizes return but violates torque/velocity limits.",
            "ρ_thetadot and ρ_torque subplots show robustness crossing 0 when constraints hold.",
        ],
        "images": [
            ("pendulum_reward.png", "Return over training: Our (brown) surpasses reward shaping (blue) after ~32 episodes."),
            ("pendulum_cost.png", "Total STL cost: Our rises toward 0 (satisfied); shaping stays near −180 (violated)."),
            ("pendulum_theta.png", "Pendulum angle |θ|: both methods learn to swing up toward upright (θ→0)."),
            ("pendulum_thetadot_subplots.png", "Angular velocity |θ̇| per agent — constraint |θ̇|≤0.5 tracked over time."),
            ("pendulum_torque_subplots.png", "Applied torque |τ| — constraint |τ|≤0.3; Our stays within bounds better."),
            ("pendulum_rho_thetadot_subplots.png", "STL robustness ρ for velocity: positive ρ ⇒ constraint satisfied."),
            ("pendulum_rho_torque_subplots.png", "STL robustness ρ for torque: Our reaches positive robustness; shaping does not."),
        ],
    },
}


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8") if path.is_file() else ""


def _style_title(title_shape, text: str, size: int = 32) -> None:
    title_shape.text = text
    p = title_shape.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(12, 44, 92)


def _add_bullet_slide(prs: Presentation, title: str, sections: list[tuple[str, list[str]]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _style_title(slide.shapes.title, title)
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    first = True
    for header, items in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = header
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(12, 44, 92)
        for item in items:
            c = tf.add_paragraph()
            c.text = item
            c.level = 1
            c.font.size = Pt(14)


def _add_image_slide(
    prs: Presentation,
    lab_num: int,
    lab_title: str,
    image_path: Path,
    caption: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _style_title(slide.shapes.title, f"Lab {lab_num} — {image_path.stem.replace('_', ' ').title()}", size=28)

    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.9))
    cap_tf = cap_box.text_frame
    cap_tf.word_wrap = True
    cap_tf.text = caption
    for p in cap_tf.paragraphs:
        p.font.size = Pt(13)
        p.font.italic = True

    slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.95), width=Inches(12.1))


def _add_csv_slide(prs: Presentation, lab_num: int, lab_title: str, csv_path: Path) -> None:
    rows: list[list[str]] = []
    with csv_path.open(newline="", encoding="utf-8") as f:
        reader = csv.reader(f)
        for idx, row in enumerate(reader):
            rows.append([str(x) for x in row])
            if idx >= 10:
                break
    if not rows:
        return

    max_cols = min(max(len(r) for r in rows), 6)
    trimmed = [(r + [""] * max_cols)[:max_cols] for r in rows]

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _style_title(slide.shapes.title, f"Lab {lab_num} — Metrics Table", size=28)
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(12.0), Inches(0.35))
    sub.text_frame.text = csv_path.name

    table_shape = slide.shapes.add_table(
        len(trimmed), max_cols, Inches(0.4), Inches(1.45), Inches(12.4), Inches(5.2)
    )
    table = table_shape.table
    for r, row in enumerate(trimmed):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11 if r else 12)
            if r == 0:
                p.font.bold = True


def _add_text_artifact_slide(
    prs: Presentation, lab_num: int, path: Path, header: str, max_lines: int = 10
) -> None:
    if not path.is_file():
        return
    lines = path.read_text(encoding="utf-8").strip().splitlines()[:max_lines]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _style_title(slide.shapes.title, f"Lab {lab_num} — Sample Output", size=28)
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = header
    p.font.bold = True
    p.font.size = Pt(16)
    for line in lines:
        c = tf.add_paragraph()
        c.text = line
        c.font.size = Pt(12)
        c.font.name = "Consolas"


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _style_title(slide.shapes.title, "Explainable AI — Labs 1–8 Portfolio")
    sub = slide.placeholders[1]
    sub.text = (
        "What was required • How we implemented it • What the results mean\n"
        "Causal XAI • Symbolic planning • Neuro-symbolic RL"
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)
    _add_bullet_slide(
        prs,
        "Course Overview",
        [
            ("Scope:", ["8 labs completed — all artifacts in assignments/"]),
            ("Structure per lab:", ["Asked → Implementation → Results → Interpretation"]),
            ("Evidence:", ["Plots, CSV tables, plans, and metrics.json"]),
        ],
    )

    for num, folder, title in LABS:
        base = ASSIGNMENTS / folder
        results = base / "results"
        narrative = LAB_NARRATIVES.get(num, {})

        _add_bullet_slide(
            prs,
            f"Lab {num} — What Was Asked",
            [("Objective:", narrative.get("asked", [title]))],
        )
        _add_bullet_slide(
            prs,
            f"Lab {num} — How We Did It",
            [("Implementation:", narrative.get("implementation", ["See SUMMARY.md"]))],
        )
        _add_bullet_slide(
            prs,
            f"Lab {num} — Results & What They Mean",
            [
                ("Measured outcomes:", narrative.get("results", [])),
                ("Interpretation:", narrative.get("interpretation", [])),
            ],
        )

        for fname, caption in narrative.get("images", []):
            img_path = results / fname
            if img_path.is_file():
                _add_image_slide(prs, num, title, img_path, caption)

        csv_name = narrative.get("csv")
        if csv_name:
            csv_path = results / csv_name
            if not csv_path.is_file():
                csvs = sorted(results.glob("*.csv"))
                csv_path = csvs[0] if csvs else None
            if csv_path and csv_path.is_file():
                _add_csv_slide(prs, num, title, csv_path)

        artifact = narrative.get("text_artifact")
        if artifact:
            fname, header = artifact
            _add_text_artifact_slide(prs, num, results / fname, header)

    _add_bullet_slide(
        prs,
        "Summary",
        [
            ("Completed:", ["All 8 assignments with reproducible results folders"]),
            ("Key theme:", ["Explainability via causality, logic, and constrained RL"]),
            ("Artifacts:", ["assignments/portfolio/ — regenerate with build_presentation.py"]),
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved presentation: {OUT}")
    slide_count = len(prs.slides)
    print(f"Total slides: {slide_count}")


if __name__ == "__main__":
    main()
